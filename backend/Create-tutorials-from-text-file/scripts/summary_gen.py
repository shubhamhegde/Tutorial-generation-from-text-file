from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from nltk.tokenize import word_tokenize, sent_tokenize
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import nltk
from pptx.util import Pt
import re,os

import create_video as cv

def processing(text):


  '''CREATING FREQUENCY TABLE'''
  stopWords = set(stopwords.words("english"))

  #print(stopWords)

  import nltk
  nltk.download('punkt')
  words = word_tokenize(text)

  #print(words)

  ps = PorterStemmer()

  freqTable = dict()
  for word in words:
      word = ps.stem(word)
      if word in stopWords:
              continue
      if word in freqTable:
              freqTable[word] += 1
      else:
              freqTable[word] = 1

  #freqTable

  '''SCORING SENTENCES'''
  sentenceValue = dict()

  sentences = text.split('.')

  for i in range(len(sentences)):
    sentences[i]=sentences[i]+"."

  #sentences

  #len(sentences)

  for sentence in sentences:
          word_count_in_sentence = (len(word_tokenize(sentence)))
          word_count_in_sentence_except_stop_words = 0
          for wordValue in freqTable:
              if wordValue in sentence.lower():
                  word_count_in_sentence_except_stop_words += 1
                  if sentence[:10] in sentenceValue:
                      sentenceValue[sentence[:10]] += freqTable[wordValue]
                  else:
                      sentenceValue[sentence[:10]] = freqTable[wordValue]

  if sentence[:10] in sentenceValue:
          sentenceValue[sentence[:10]] = sentenceValue[sentence[:10]]

  #sentenceValue

  '''FINDING AVERAGE SCORE'''
  sumValues = 0
  for entry in sentenceValue:
      sumValues += sentenceValue[entry]

      # Average value of a sentence from original text
      average = (sumValues / len(sentenceValue))

  #average

  sentence_count = 0
  summary = ''
  global mapping
  mapping={}
  start=0
  for i in range(len(sentences)):
      sentence=sentences[i]
      if sentence[:10] in sentenceValue and sentenceValue[sentence[:10]] >= (average):
              summary += " " + sentence
              sentence_count += 1
              mapping[sentence_count-1]=i
              #start=i+1
  
  print(mapping)

  print(summary)
  return summary,mapping


def pptgen(text,summary,mapping,filename):

  language = 'hi'
  sentences = text.split('.')

  for i in range(len(sentences)):
    sentences[i]=sentences[i]+"."

  from gtts import gTTS 

  summary = summary.split('.')
  summary.pop()

  for i in range(len(summary)):
    if summary[i][1]=="\n":
      summary[i]=summary[i][2:]+"."
    else:
      summary[i]=summary[i][1:]+"."

  print("xxxxxxxxxxxx")
  print(summary)
  print("xxxxxxxxxxxx")

  #summary

  from pptx import Presentation,util

  prs = Presentation()
  title_slide_layout = prs.slide_layouts[1]
  slide_to_voice={}
  #Used to add bullet points
  #tf = body_shape.text_frame
  for i in range(0,len(summary),3):
    slide = prs.slides.add_slide(title_slide_layout)
    print(slide.element)
    title = slide.shapes.title
    if i==0:
      title.text="Abstract"
    subtitle = slide.placeholders[1]
    left = top = width = height = util.Inches(1.0)
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    # tf.fit_text()# = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
    title_shape = shapes.title
    x=min(i+3,len(summary))
    for j in range(i,x):
      p = tf.add_paragraph()
      p.text = summary[j]
      p.level=0
    tf.fit_text(max_size=25)
    print("^^^^^^^^^^^^^^^^")
    print(tf.paragraphs[0].font.size)
    print("^^^^^^^^^^^^^^^^")
    # for paragraph in tf.paragraphs:
    #   for run in paragraph.runs:
    #       run.font.size = Pt(9)
    voice_text=' '.join(sentences[mapping[str(i)]:mapping[str(x-1)]+1])
    myobj = gTTS(text=voice_text, lang=language, slow=False) 
    myobj.save("voiceover_"+str(i)+".mp3")
    slide_to_voice[i]="voiceover_"+str(i)+".mp3"
    movie = slide.shapes.add_movie("voiceover_"+str(i)+".mp3", 
                    left , top , width , height, 
                              poster_frame_image=None, 
                              mime_type='video/unknown')
  ppt_path='static/downloads/'+filename
  prs.save(ppt_path)
  name=filename.split(".")[0]
  pdf_name=name+".pdf"
  video_name=name+".mp4"
  ppt_path=os.path.abspath(ppt_path)
  pdf_path=ppt_path[:-4]+"pdf"
  print("**********    ",pdf_path,"      *********")
  #cv.ppt_presenter(filename,pdf_name,video_name,slide_to_voice)
  #cv.ppt_presenter('D:\College\Capstone project\Create-tutorials-from-text-file\scripts\static\downloads\hi_summary.pptx','D:\College\Capstone project\Create-tutorials-from-text-file\scripts\static\downloads\hi_summary.pdf','D:\College\Capstone project\Create-tutorials-from-text-file\scripts\static\downloads\hi_summary.mp4',slide_to_voice)
  cv.PPTtoPDF(ppt_path,pdf_path)
  return {'ppt_path':ppt_path,'pdf_path':pdf_path}


