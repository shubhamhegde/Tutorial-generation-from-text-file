import os
import tempfile
import argparse
from subprocess import call

from pdf2image import convert_from_path
from pptx import Presentation
from gtts import gTTS
import pythoncom

FFMPEG_NAME = r'C:\Users\Kritika Kapoor\Documents\ShareX\Tools\ffmpeg.exe'

#import comtypes.client
import win32com.client

def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    #CoInitialize()
    #powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    pythoncom.CoInitialize()
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName,WithWindow=False)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()

def ffmpeg_call(image_path, audio_path, temp_path, i):
    out_path_mp4 = os.path.join(temp_path, 'frame_{}.mp4'.format(i))
    out_path_ts = os.path.join(temp_path, 'frame_{}.ts'.format(i))
    print(out_path_mp4)
    call([FFMPEG_NAME, '-loop', '1', '-y', '-i', image_path, '-i', audio_path,'-c:v', 'libx264', '-tune', 'stillimage', '-c:a', 'aac','-b:a', '192k', '-pix_fmt', 'yuv420p', '-shortest',out_path_mp4],shell=True)
    call([FFMPEG_NAME, '-y', '-i', out_path_mp4, '-c', 'copy',
          '-bsf:v', 'h264_mp4toannexb', '-f', 'mpegts', out_path_ts],shell=True)

def ffmpeg_concat(video_list_str, out_path):
    call([FFMPEG_NAME, '-y', '-f', 'mpegts', '-i', '{}'.format(video_list_str),
          '-c', 'copy', '-bsf:a', 'aac_adtstoasc', out_path])

def ppt_presenter(pptx_path, pdf_path, output_path, voiceover):
    PPTtoPDF(pptx_path,pdf_path)
    with tempfile.TemporaryDirectory() as temp_path:
        images_from_path = convert_from_path(pdf_path,poppler_path=r'C:\Users\Kritika Kapoor\Downloads\poppler-0.68.0_x86\poppler-0.68.0\bin')
        prs = Presentation(pptx_path)
        assert len(images_from_path) == len(prs.slides)
        for i, (slide, image) in enumerate(zip(prs.slides, images_from_path)):
            #if slide.has_notes_slide:
            #notes = slide.notes_slide.notes_text_frame.text
            #tts = gTTS(text=notes, lang='en')
            image_path = os.path.join(temp_path, 'frame_{}.jpg'.format(i))
            audio_path = os.path.abspath(voiceover[i*3])

            image.save(image_path)
            #tts.save(audio_path)
            print(image_path,audio_path,temp_path)
            ffmpeg_call(image_path, audio_path, temp_path, i)

        video_list = [os.path.join(temp_path, 'frame_{}.ts'.format(i)) \
                      for i in range(len(images_from_path))]
        video_list_str = 'concat:' + '|'.join(video_list)
        ffmpeg_concat(video_list_str, output_path)


