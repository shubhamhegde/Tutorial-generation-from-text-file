from operator import itemgetter
import fitz
import json
import summary_gen as sg
from pptx.util import Pt
import os

def fonts(doc, granularity=False):
    """Extracts fonts and their usage in PDF documents.
    :param doc: PDF document to iterate through
    :type doc: <class 'fitz.fitz.Document'>
    :param granularity: also use 'font', 'flags' and 'color' to discriminate text
    :type granularity: bool
    :rtype: [(font_size, count), (font_size, count}], dict
    :return: most used fonts sorted by count, font style information
    """
    styles = {}
    font_counts = {}

    for page in doc:
        blocks = page.getText("dict")["blocks"]
        for b in blocks:  # iterate through the text blocks
            if b['type'] == 0:  # block contains text
                for l in b["lines"]:  # iterate through the text lines
                    for s in l["spans"]:  # iterate through the text spans
                        if granularity:
                            identifier = "{0}_{1}_{2}_{3}".format(s['size'], s['flags'], s['font'], s['color'])
                            styles[identifier] = {'size': s['size'], 'flags': s['flags'], 'font': s['font'],
                                                  'color': s['color']}
                        else:
                            identifier = "{0}".format(s['size'])
                            styles[identifier] = {'size': s['size'], 'font': s['font']}

                        font_counts[identifier] = font_counts.get(identifier, 0) + 1  # count the fonts usage

    font_counts = sorted(font_counts.items(), key=itemgetter(1), reverse=True)

    if len(font_counts) < 1:
        raise ValueError("Zero discriminating fonts found!")

    return font_counts, styles


def font_tags(font_counts, styles):
    """Returns dictionary with font sizes as keys and tags as value.
    :param font_counts: (font_size, count) for all fonts occuring in document
    :type font_counts: list
    :param styles: all styles found in the document
    :type styles: dict
    :rtype: dict
    :return: all element tags based on font-sizes
    """
    p_style = styles[font_counts[0][0]]  # get style for most used font by count (paragraph)
    p_size = p_style['size']  # get the paragraph's size

    # sorting the font sizes high to low, so that we can append the right integer to each tag
    font_sizes = []
    for (font_size, count) in font_counts:
        font_sizes.append(float(font_size))
    font_sizes.sort(reverse=True)

    # aggregating the tags for each font size
    idx = 0
    size_tag = {}
    for size in font_sizes:
        idx += 1
        if size == p_size:
            idx = 0
            size_tag[size] = '<p>'
        if size > p_size:
            size_tag[size] = '<h{0}>'.format(idx)
        elif size < p_size:
            size_tag[size] = '<s{0}>'.format(idx)

    return size_tag


def headers_para(doc, size_tag):
    """Scrapes headers & paragraphs from PDF and return texts with element tags.
    :param doc: PDF document to iterate through
    :type doc: <class 'fitz.fitz.Document'>
    :param size_tag: textual element tags for each size
    :type size_tag: dict
    :rtype: list
    :return: texts with pre-prended element tags
    """
    header_para = []  # list with headers and paragraphs
    first = True  # boolean operator for first header
    previous_s = {}  # previous span
    block_string = ""
    for page in doc:
        blocks = page.getText("dict")["blocks"]
        #print("Blocks --------- ",blocks)
        for b in blocks:  # iterate through the text blocks
            if b['type'] == 0:  # this block contains text

                # REMEMBER: multiple fonts and sizes are possible IN one block

                #block_string = ""  # text found in block
                for l in b["lines"]:  # iterate through the text lines
                    for s in l["spans"]:  # iterate through the text spans
                        #print(s)
                        if len(s['text'])>1 and s['text'].strip():  # removing whitespaces:
                            s['text']=s['text'].encode('ascii','ignore').decode('utf-8')
                            if first:
                                previous_s = s
                                first = False
                                block_string = size_tag[s['size']] + s['text']
                            else:
                                if s['size'] == previous_s['size']:
                                    #if block_string: #and all((c == "|") for c in block_string):
                                    # block_string only contains pipes
                                    #    block_string = size_tag[s['size']] + s['text']
                                    #if block_string == "":
                                    #    # new block has started, so append size tag
                                    #    block_string = size_tag[s['size']] + s['text']
                                    #else:  # in the same block, so concatenate strings
                                    block_string += " " + s['text']

                                else:
                                    header_para.append(block_string)
                                    block_string = size_tag[s['size']] + s['text']

                                previous_s = s
                    #print("-------------")

                    # new block started, indicating with a pipe
                    #block_string += "|"

                #header_para.append(block_string)
                #print("**************")
    header_para.append(block_string)
    return header_para


def preprocessing(filename):

    document = filename
    doc = fitz.open(document)

    font_counts, styles = fonts(doc, granularity=False)
    #print(font_counts,styles)
    size_tag = font_tags(font_counts, styles)
    #print(size_tag)
    elements = headers_para(doc, size_tag)
    #print(elements)
    with open("doc.json", 'w') as json_out:
        json.dump(elements, json_out)
    new_elements=[]
    prev_tag=""
    return elements

def remove_string_special_characters(s): 
        
        # removes special characters with ' ' 
        stripped = re.sub('[^a-zA-z\s]', '', s) 
        stripped = re.sub('_', '', stripped) 
        
        # Change any white space to one space 
        stripped = re.sub('\s+', ' ', stripped) 
        
        # Remove start and end white spaces 
        stripped = stripped.strip() 
        if stripped != '': 
                return stripped.lower()

def topic_gen(text):
    import nltk 
    import re 
    from sklearn.feature_extraction.text import CountVectorizer, TfidfVectorizer 
    from nltk.corpus import stopwords 
    from nltk.tokenize import word_tokenize 
    import pandas as pd 

    
    txt1=text.split(".")

    # Preprocessing 
            
    # Stopword removal 
    stop_words = set(stopwords.words('english')) 
    your_list = ['skills', 'ability', 'job', 'description'] 
    for i, line in enumerate(txt1): 
        txt1[i] = ' '.join([x for
            x in nltk.word_tokenize(line) if
            ( x not in stop_words ) and ( x not in your_list )])
        
    # Getting trigrams 
    vectorizer = CountVectorizer(ngram_range = (3,3)) 
    X1 = vectorizer.fit_transform(txt1) 
    features = (vectorizer.get_feature_names()) 
    # print("\n\nFeatures : \n", features) 
    # print("\n\nX1 : \n", X1.toarray()) 

    # Applying TFIDF 
    vectorizer = TfidfVectorizer(ngram_range = (3,3)) 
    X2 = vectorizer.fit_transform(txt1) 
    scores = (X2.toarray()) 
    # print("\n\nScores : \n", scores) 

    # Getting top ranking features 
    sums = X2.sum(axis = 0) 
    data1 = []
    features=list(features)
    for col, term in enumerate(features): 
        data1.append( (term, sums[0,col] )) 
    ranking = pd.DataFrame(data1, columns = ['term','rank']) 
    words3 = (ranking.sort_values('rank', ascending = False)).reset_index()
    print ("\n\nWords head : \n", words3.head(7))

    vectorizer = CountVectorizer(ngram_range = (2,2)) 
    X1 = vectorizer.fit_transform(txt1) 
    features = (vectorizer.get_feature_names()) 
    # print("\n\nFeatures : \n", features) 
    # print("\n\nX1 : \n", X1.toarray()) 

    # Applying TFIDF 
    vectorizer = TfidfVectorizer(ngram_range = (2,2)) 
    X2 = vectorizer.fit_transform(txt1) 
    scores = (X2.toarray()) 
    # print("\n\nScores : \n", scores) 

    # Getting top ranking features 
    sums = X2.sum(axis = 0) 
    data1 = []
    features=list(features)
    for col, term in enumerate(features): 
        data1.append( (term, sums[0,col] )) 
    ranking = pd.DataFrame(data1, columns = ['term','rank']) 
    words2 = (ranking.sort_values('rank', ascending = False)).reset_index()
    print ("\n\nWords head : \n", words2.head(7))

    vectorizer = CountVectorizer(ngram_range = (1,1)) 
    X1 = vectorizer.fit_transform(txt1) 
    features = (vectorizer.get_feature_names()) 
    # print("\n\nFeatures : \n", features) 
    # print("\n\nX1 : \n", X1.toarray()) 

    # Applying TFIDF 
    vectorizer = TfidfVectorizer(ngram_range = (1,1)) 
    X2 = vectorizer.fit_transform(txt1) 
    scores = (X2.toarray()) 
    # print("\n\nScores : \n", scores) 

    # Getting top ranking features 
    sums = X2.sum(axis = 0) 
    data1 = []
    features=list(features)
    for col, term in enumerate(features): 
        data1.append( (term, sums[0,col] )) 
    ranking = pd.DataFrame(data1, columns = ['term','rank']) 
    words1 = (ranking.sort_values('rank', ascending = False)).reset_index()
    print ("\n\nWords head : \n", words1.head(7))

    threshold=(words1['rank'][0]+words2['rank'][0]+words3['rank'][0])/3
    print(words1['term'].loc[0],threshold)
    if words2['rank'][0]>=threshold:
        return words2['term'][0]
    elif words3['rank'][0]>=threshold:
        return words3['term'][0]
    else:
        return words1['term'][0]
    

def pptgen(elements,filename,tid):

  
  from pptx import Presentation,util

  prs = Presentation()
  title_slide_layout = prs.slide_layouts[1]
  title_layout = prs.slide_layouts[0]
  slide_to_voice={}
  subtopic_mapping=[]
  #Used to add bullet points
  #tf = body_shape.text_frame
  pg_cnt=0
  if len(elements)>1:
    current=0
    next_element=1
    while(next_element<len(elements)):
        if elements[next_element][:3]=="<p>":
                pg_cnt+=1
                heading=int(elements[current][2])
                if heading==1:
                    subtopic_mapping.append({"heading":elements[current][4:],"level":1,"pgno":pg_cnt,"children":[]})
                else:
                    x=subtopic_mapping[-1]
                    while(len(x)>=1):
                        if x["level"]==heading-1:
                            x["children"].append({"heading":elements[current][4:],"level":heading,"pgno":pg_cnt,"children":[]})
                            break
                        else:
                            if not x["children"] or x["children"][-1]["level"]==heading:
                                x["children"].append({"heading":elements[current][4:],"level":heading,"pgno":pg_cnt,"children":[]})
                                break
                            else:
                                x=x["children"][-1]
                summary=sg.processing(elements[next_element][3:])[0]
                summary = summary.split('.')
                summary.pop()

                for i in range(len(summary)):
                    if summary[i][1]=="\n":
                        summary[i]=summary[i][2:]+"."
                    else:
                        summary[i]=summary[i][1:]+"."
                for i in range(0,len(summary),3):
                    slide = prs.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    if i==0:
                        title.text=elements[current][4:]
                    else:
                        pg_cnt+=1
                        title.text=elements[current][4:]+" contd."
                    subtitle = slide.placeholders[1]
                    left = top = width = height = util.Inches(1.0)
                    shapes = slide.shapes
                    body_shape = shapes.placeholders[1]
                    tf = body_shape.text_frame
                    # tf.fit_text()# = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    title_shape = shapes.title
                    x=min(i+3,len(summary))
                    for j in range(i,x):
                        p = tf.add_paragraph()
                        p.text = summary[j]
                        p.level=0
                    tf.fit_text(max_size=25)
                    #print("^^^^^^^^^^^^^^^^")
                    #print(tf.paragraphs[0].font.size)
                    #print("^^^^^^^^^^^^^^^^")
        else:
            if elements[current][:2]=="<h":
                pg_cnt+=1
                heading=int(elements[current][2])
                if heading==1:
                    subtopic_mapping.append({"heading":elements[current][4:],"level":1,"pgno":pg_cnt,"children":[]})
                else:
                    x=subtopic_mapping[-1]
                    while(len(x)>=1):
                        if x["level"]==heading-1:
                            x["children"].append({"heading":elements[current][4:],"level":heading,"pgno":pg_cnt,"children":[]})
                            break
                        else:
                            x=x["children"][-1]
                slide = prs.slides.add_slide(title_layout)
                title = slide.shapes.title
                title.text = elements[current][4:]
                title.text_frame.paragraphs[0].font.size=Pt(60)
        current+=1
        next_element+=1
  
  else:
        print("here")
        title_text=topic_gen(elements[0][3:])
        subtopic_mapping.append({"heading":title_text.title(),"level":1,"pgno":1,"children":[]})
        slide = prs.slides.add_slide(title_layout)
        title = slide.shapes.title
        title.text = title_text.upper()
        title.text_frame.paragraphs[0].font.size=Pt(60)
        summary=sg.processing(elements[0][3:])[0]
        summary = summary.split('.')
        summary.pop()

        for i in range(len(summary)):
            if summary[i][1]=="\n":
                summary[i]=summary[i][2:]+"."
            else:
                summary[i]=summary[i][1:]+"."
        for i in range(0,len(summary),3):
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            if i==0:
                title.text=title_text.title()
            else:
                title.text=title_text.title()+" contd."
            left = top = width = height = util.Inches(1.0)
            shapes = slide.shapes
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            # tf.fit_text()# = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            title_shape = shapes.title
            x=min(i+3,len(summary))
            for j in range(i,x):
                p = tf.add_paragraph()
                p.text = summary[j]
                p.level=0
            tf.fit_text(max_size=25)
                
  ppt_path='static/downloads/'+tid+".pptx"
  prs.save(ppt_path)
  print("**************")
  print(subtopic_mapping)
  name=filename.split(".")[0]
  pdf_name=name+".pdf"
  video_name=name+".mp4"
  ppt_path=os.path.abspath(ppt_path)
  pdf_path=ppt_path[:-4]+"pdf"
  print("**********    ",pdf_path,"      *********")
  #cv.ppt_presenter(filename,pdf_name,video_name,slide_to_voice)
  #cv.ppt_presenter('D:\College\Capstone project\Create-tutorials-from-text-file\scripts\static\downloads\hi_summary.pptx','D:\College\Capstone project\Create-tutorials-from-text-file\scripts\static\downloads\hi_summary.pdf','D:\College\Capstone project\Create-tutorials-from-text-file\scripts\static\downloads\hi_summary.mp4',slide_to_voice)
  #cv.PPTtoPDF(ppt_path,pdf_path)
  return {'ppt_path':ppt_path,'pdf_path':pdf_path,'mapping':subtopic_mapping}
